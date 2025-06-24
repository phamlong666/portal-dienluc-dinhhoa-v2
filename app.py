import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
import io
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptInches

st.set_page_config(page_title="Báo cáo tổn thất TBA", layout="wide")
st.title("📊 Báo cáo tổn thất các TBA công cộng")

file_keys = ["Theo Tháng", "Lũy kế", "Cùng kỳ"]
uploaded_data = {}

col_uploads = st.columns(3)
for i, key in enumerate(file_keys):
    with col_uploads[i]:
        file = st.file_uploader(f"📁 File {key}", type=["xlsx"], key=f"upload_{key}")
        if file:
            xls = pd.ExcelFile(file)
            sheet_name = [s for s in xls.sheet_names if "ánh xạ" in s.lower()][0]
            df = pd.read_excel(xls, sheet_name=sheet_name)
            uploaded_data[key] = df

if uploaded_data:
    if st.button("📌 Tạo báo cáo"):

        def export_powerpoint(title, actual, plan):
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = title

            txBox = shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(5))
            tf = txBox.text_frame
            tf.text = f"Tỷ lệ tổn thất thực tế: {actual:.2f}%"
            p = tf.add_paragraph()
            p.text = f"Tỷ lệ tổn thất kế hoạch: {plan:.2f}%"
            return prs

        for key, df in uploaded_data.items():
            st.subheader(f"🔍 Dữ liệu: {key}")
            df_copy = df.copy()
            percent_cols = [col for col in df_copy.columns if "%" in col]
            for col in percent_cols:
                df_copy[col] = pd.to_numeric(df_copy[col], errors="coerce")
                df_copy[col] = df_copy[col].map(lambda x: f"{x:.2f}%" if pd.notna(x) else "")
            st.dataframe(df_copy.style.set_properties(**{"font-size": "18pt"}), use_container_width=True)

            total_input = df["Điện nhận (kWh)"].sum()
            total_loss = df["Điện tổn thất (kWh)"].sum()
            actual = (total_loss / total_input * 100) if total_input else 0
            plan_col = [c for c in df.columns if "kế hoạch" in c.lower()][0]
            plan_series = df[plan_col]
            plan = ((plan_series / 100 * df["Điện nhận (kWh)"]).sum() / total_input * 100) if total_input else 0

            st.markdown(f"#### 📉 Biểu đồ tổn thất - {key}")
            fig = go.Figure(data=[
                go.Bar(
                    x=["Thực tế", "Kế hoạch"],
                    y=[actual, plan],
                    marker=dict(
                        color=["#3366cc", "#ff9900"],
                        line=dict(color='black', width=1)
                    ),
                    text=[f"{actual:.2f}%", f"{plan:.2f}%"],
                    textposition='auto',
                    width=[0.4, 0.4]
                )
            ])
            fig.update_layout(
                title="Tỷ lệ tổn thất điện năng",
                margin=dict(l=20, r=20, t=40, b=20),
                height=350,
                font=dict(size=16),
                showlegend=False,
                yaxis=dict(title="Tỷ lệ (%)", range=[0, max(actual, plan) * 1.2 if max(actual, plan) > 0 else 5])
            )
            st.plotly_chart(fig, use_container_width=True)

        if len(uploaded_data) == 3:
            st.markdown("### 📊 Biểu đồ hợp nhất tổn thất các file")
            data_total = []
            for key in file_keys:
                df = uploaded_data[key]
                total_input = df["Điện nhận (kWh)"].sum()
                total_loss = df["Điện tổn thất (kWh)"].sum()
                actual = (total_loss / total_input * 100) if total_input else 0
                plan_col = [c for c in df.columns if "kế hoạch" in c.lower()][0]
                plan_series = df[plan_col]
                plan = ((plan_series / 100 * df["Điện nhận (kWh)"]).sum() / total_input * 100) if total_input else 0
                data_total.append((key, actual, plan))

            fig2 = go.Figure()
            x = file_keys
            actuals = [d[1] for d in data_total]
            plans = [d[2] for d in data_total]
            fig2.add_trace(go.Bar(
                x=x,
                y=actuals,
                name='Thực tế',
                marker_color='#3366cc',
                text=[f"{v:.2f}%" for v in actuals],
                textposition='auto'
            ))
            fig2.add_trace(go.Bar(
                x=x,
                y=plans,
                name='Kế hoạch',
                marker_color='#ff9900',
                text=[f"{v:.2f}%" for v in plans],
                textposition='auto'
            ))
            fig2.update_layout(
                barmode='group',
                height=350,
                margin=dict(l=30, r=30, t=30, b=30),
                font=dict(size=16),
                yaxis=dict(title="Tỷ lệ (%)", range=[0, max(actuals + plans) * 1.2 if actuals else 5])
            )
            st.plotly_chart(fig2, use_container_width=True)
# ============== Biểu đồ tổn thất theo ngưỡng =====================

import matplotlib.pyplot as plt

def ve_bieu_do_nguong_ton_that():
    categories = ["<2%", ">=2 và <3%", ">=3 và <4%", ">=4 và <5%", ">=5 và <7%", ">=7%"]
    values_cungky = [4, 2, 24, 42, 86, 35]
    values_thuchien = [4, 14, 32, 55, 69, 29]

    fig1, ax = plt.subplots(figsize=(10, 5))
    x = range(len(categories))
    ax.bar([i - 0.2 for i in x], values_cungky, width=0.4, label="Cùng kỳ", color="lightgray")
    ax.bar([i + 0.2 for i in x], values_thuchien, width=0.4, label="Thực hiện", color="teal")
    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=14)
    ax.set_ylabel("Số lượng TBA", fontsize=16)
    ax.set_title("Số lượng TBA theo ngưỡng tổn thất", fontsize=18)
    ax.legend()
    ax.grid(axis='y', linestyle='--', alpha=0.5)
    st.pyplot(fig1)

    # Donut chart
    sizes = [6.90, 6.90, 15.76, 27.09, 33.99, 14.29]
    colors = ["#1f77b4", "#ff9900", "#2ca02c", "#bcbd22", "#17becf", "#d62728"]
    fig2, ax2 = plt.subplots(figsize=(6, 6))
    wedges, texts, autotexts = ax2.pie(sizes, labels=categories, colors=colors,
                                       autopct='%1.1f%%', startangle=90,
                                       wedgeprops=dict(width=0.3))
    ax2.set(aspect="equal", title='Tỷ trọng TBA theo ngưỡng tổn thất')
    st.pyplot(fig2)

if st.button("📊 Biểu đồ theo ngưỡng tổn thất"):
    ve_bieu_do_nguong_ton_that()
