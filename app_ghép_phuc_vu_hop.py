
import streamlit as st

# ================== Giao diện Điều hành số ==================
tab = st.session_state.get("tab", "Trang chính")

if tab == "Trang chính":
    # Nút chuyển sang phục vụ họp
    if st.button("🧾 Phục vụ họp – Ghi báo cáo và xuất file"):
        st.session_state["tab"] = "Phục vụ họp"
        st.experimental_rerun()
    # Hiển thị giao diện chính (các nút lớn khác giữ nguyên theo original_app)
    exec("""
import streamlit as st
import datetime
from docx import Document
from pptx import Presentation
from io import BytesIO
import csv

# Hàm lưu vào CSV
def save_to_csv(ten, ngay, noidung):
    csv_file = "lich_su_cuoc_hop.csv"
    header = ["ten_cuoc_hop", "ngay_hop", "noi_dung", "ngay_ghi"]
    try:
        with open(csv_file, mode="a", encoding="utf-8", newline="") as file:
            writer = csv.DictWriter(file, fieldnames=header)
            if file.tell() == 0:
                writer.writeheader()
            writer.writerow({
                "ten_cuoc_hop": ten,
                "ngay_hop": ngay.strftime("%d/%m/%Y"),
                "noi_dung": noidung,
                "ngay_ghi": datetime.datetime.now().strftime("%d/%m/%Y %H:%M")
            })
        return csv_file
    except Exception as e:
        return str(e)

# Hàm tạo Word
def tao_bao_cao_word(ten, ngay, noidung):
    doc = Document()
    doc.add_heading("BÁO CÁO CUỘC HỌP", 0)
    doc.add_paragraph(f"Tên cuộc họp: {ten}")
    doc.add_paragraph(f"Ngày họp: {ngay.strftime('%d/%m/%Y')}")
    doc.add_paragraph("\nNội dung chính:")
    doc.add_paragraph(noidung)
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# Hàm tạo PowerPoint
def tao_bao_cao_ppt(ten, ngay, noidung):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "BÁO CÁO CUỘC HỌP"
    slide.placeholders[1].text = f"{ten}\nNgày họp: {ngay.strftime('%d/%m/%Y')}"
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide.shapes.title.text = "Nội dung chính"
    content_slide.placeholders[1].text = noidung
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Giao diện Streamlit
st.title("📋 Phục vụ họp – Ghi báo cáo và xuất file")

cuoc_hop = st.text_input("Tên cuộc họp")
ngay_hop = st.date_input("Ngày họp", value=datetime.date.today())
noi_dung = st.text_area("Nội dung cuộc họp", height=300)

col1, col2, col3 = st.columns(3)
with col1:
    if st.button("📤 Tạo Word"):
        if cuoc_hop and noi_dung:
            word_file = tao_bao_cao_word(cuoc_hop, ngay_hop, noi_dung)
            st.download_button("📥 Tải Word", word_file, file_name=f"{cuoc_hop.replace(' ','_')}.docx")
            save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        else:
            st.warning("❗ Nhập tên và nội dung họp.")
with col2:
    if st.button("📽️ Tạo PowerPoint"):
        if cuoc_hop and noi_dung:
            ppt_file = tao_bao_cao_ppt(cuoc_hop, ngay_hop, noi_dung)
            st.download_button("📥 Tải PowerPoint", ppt_file, file_name=f"{cuoc_hop.replace(' ','_')}.pptx")
            save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        else:
            st.warning("❗ Nhập tên và nội dung họp.")
with col3:
    if st.button("📜 Lưu lịch sử"):
        path = save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        st.success(f"✅ Đã lưu vào {path}")

# Hiển thị lịch sử nếu có
st.markdown("---")
st.subheader("📚 Lịch sử cuộc họp đã lưu")
try:
    with open("lich_su_cuoc_hop.csv", encoding="utf-8") as f:
        reader = csv.reader(f)
        rows = list(reader)
        for row in rows[1:]:
            st.markdown(f"**🗓 {row[1]}** – `{row[0]}`  
{row[2]}")
except:
    st.info("Chưa có cuộc họp nào được lưu.")
""")

elif tab == "Phục vụ họp":
    st.markdown("### 🧾 Phục vụ họp – Ghi báo cáo và xuất file")
    if st.button("🔙 Quay về trang chính"):
        st.session_state["tab"] = "Trang chính"
        st.experimental_rerun()
    # Giao diện nhập nội dung họp
    exec("""
import streamlit as st
import datetime
from docx import Document
from pptx import Presentation
from io import BytesIO
import csv

# Hàm lưu vào CSV
def save_to_csv(ten, ngay, noidung):
    csv_file = "lich_su_cuoc_hop.csv"
    header = ["ten_cuoc_hop", "ngay_hop", "noi_dung", "ngay_ghi"]
    try:
        with open(csv_file, mode="a", encoding="utf-8", newline="") as file:
            writer = csv.DictWriter(file, fieldnames=header)
            if file.tell() == 0:
                writer.writeheader()
            writer.writerow({
                "ten_cuoc_hop": ten,
                "ngay_hop": ngay.strftime("%d/%m/%Y"),
                "noi_dung": noidung,
                "ngay_ghi": datetime.datetime.now().strftime("%d/%m/%Y %H:%M")
            })
        return csv_file
    except Exception as e:
        return str(e)

# Hàm tạo Word
def tao_bao_cao_word(ten, ngay, noidung):
    doc = Document()
    doc.add_heading("BÁO CÁO CUỘC HỌP", 0)
    doc.add_paragraph(f"Tên cuộc họp: {ten}")
    doc.add_paragraph(f"Ngày họp: {ngay.strftime('%d/%m/%Y')}")
    doc.add_paragraph("\nNội dung chính:")
    doc.add_paragraph(noidung)
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# Hàm tạo PowerPoint
def tao_bao_cao_ppt(ten, ngay, noidung):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "BÁO CÁO CUỘC HỌP"
    slide.placeholders[1].text = f"{ten}\nNgày họp: {ngay.strftime('%d/%m/%Y')}"
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide.shapes.title.text = "Nội dung chính"
    content_slide.placeholders[1].text = noidung
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Giao diện Streamlit
st.title("📋 Phục vụ họp – Ghi báo cáo và xuất file")

cuoc_hop = st.text_input("Tên cuộc họp")
ngay_hop = st.date_input("Ngày họp", value=datetime.date.today())
noi_dung = st.text_area("Nội dung cuộc họp", height=300)

col1, col2, col3 = st.columns(3)
with col1:
    if st.button("📤 Tạo Word"):
        if cuoc_hop and noi_dung:
            word_file = tao_bao_cao_word(cuoc_hop, ngay_hop, noi_dung)
            st.download_button("📥 Tải Word", word_file, file_name=f"{cuoc_hop.replace(' ','_')}.docx")
            save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        else:
            st.warning("❗ Nhập tên và nội dung họp.")
with col2:
    if st.button("📽️ Tạo PowerPoint"):
        if cuoc_hop and noi_dung:
            ppt_file = tao_bao_cao_ppt(cuoc_hop, ngay_hop, noi_dung)
            st.download_button("📥 Tải PowerPoint", ppt_file, file_name=f"{cuoc_hop.replace(' ','_')}.pptx")
            save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        else:
            st.warning("❗ Nhập tên và nội dung họp.")
with col3:
    if st.button("📜 Lưu lịch sử"):
        path = save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        st.success(f"✅ Đã lưu vào {path}")

# Hiển thị lịch sử nếu có
st.markdown("---")
st.subheader("📚 Lịch sử cuộc họp đã lưu")
try:
    with open("lich_su_cuoc_hop.csv", encoding="utf-8") as f:
        reader = csv.reader(f)
        rows = list(reader)
        for row in rows[1:]:
            st.markdown(f"**🗓 {row[1]}** – `{row[0]}`  
{row[2]}")
except:
    st.info("Chưa có cuộc họp nào được lưu.")
""")

