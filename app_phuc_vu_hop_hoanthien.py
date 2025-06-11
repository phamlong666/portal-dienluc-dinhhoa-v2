
import streamlit as st
import datetime
from docx import Document
from pptx import Presentation
from io import BytesIO
import csv

# HÃ m lÆ°u vÃ o CSV
def save_to_csv(ten, ngay, noidung):
    csv_file = "lich_su_cuoc_hop.csv"
    header = ["ten_cuoc_hop", "ngay_hop", "noi_dung", "ngay_ghi"]
    try:
        with open(csv_file, mode="a", encoding="utf-8", newline="") as file:
            writer = csv.DictWriter(file, fieldnames=header)
            if file.tell() == 0:
                writer.writeheader()
            writer.writerow({
                "ten_cuoc_hop": ten,
                "ngay_hop": ngay.strftime("%d/%m/%Y"),
                "noi_dung": noidung,
                "ngay_ghi": datetime.datetime.now().strftime("%d/%m/%Y %H:%M")
            })
        return csv_file
    except Exception as e:
        return str(e)

# HÃ m táº¡o Word
def tao_bao_cao_word(ten, ngay, noidung):
    doc = Document()
    doc.add_heading("BÃO CÃO CUá»˜C Há»ŒP", 0)
    doc.add_paragraph(f"TÃªn cuá»™c há»p: {ten}")
    doc.add_paragraph(f"NgÃ y há»p: {ngay.strftime('%d/%m/%Y')}")
    doc.add_paragraph("\nNá»™i dung chÃ­nh:")
    doc.add_paragraph(noidung)
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# HÃ m táº¡o PowerPoint
def tao_bao_cao_ppt(ten, ngay, noidung):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "BÃO CÃO CUá»˜C Há»ŒP"
    slide.placeholders[1].text = f"{ten}\nNgÃ y há»p: {ngay.strftime('%d/%m/%Y')}"
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide.shapes.title.text = "Ná»™i dung chÃ­nh"
    content_slide.placeholders[1].text = noidung
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Giao diá»‡n Streamlit
st.title("ğŸ“‹ Phá»¥c vá»¥ há»p â€“ Ghi bÃ¡o cÃ¡o vÃ  xuáº¥t file")

cuoc_hop = st.text_input("TÃªn cuá»™c há»p")
ngay_hop = st.date_input("NgÃ y há»p", value=datetime.date.today())
noi_dung = st.text_area("Ná»™i dung cuá»™c há»p", height=300)

col1, col2, col3 = st.columns(3)
with col1:
    if st.button("ğŸ“¤ Táº¡o Word"):
        if cuoc_hop and noi_dung:
            word_file = tao_bao_cao_word(cuoc_hop, ngay_hop, noi_dung)
            st.download_button("ğŸ“¥ Táº£i Word", word_file, file_name=f"{cuoc_hop.replace(' ','_')}.docx")
            save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        else:
            st.warning("â— Nháº­p tÃªn vÃ  ná»™i dung há»p.")
with col2:
    if st.button("ğŸ“½ï¸ Táº¡o PowerPoint"):
        if cuoc_hop and noi_dung:
            ppt_file = tao_bao_cao_ppt(cuoc_hop, ngay_hop, noi_dung)
            st.download_button("ğŸ“¥ Táº£i PowerPoint", ppt_file, file_name=f"{cuoc_hop.replace(' ','_')}.pptx")
            save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        else:
            st.warning("â— Nháº­p tÃªn vÃ  ná»™i dung há»p.")
with col3:
    if st.button("ğŸ“œ LÆ°u lá»‹ch sá»­"):
        path = save_to_csv(cuoc_hop, ngay_hop, noi_dung)
        st.success(f"âœ… ÄÃ£ lÆ°u vÃ o {path}")

# Hiá»ƒn thá»‹ lá»‹ch sá»­ náº¿u cÃ³
st.markdown("---")
st.subheader("ğŸ“š Lá»‹ch sá»­ cuá»™c há»p Ä‘Ã£ lÆ°u")
try:
    with open("lich_su_cuoc_hop.csv", encoding="utf-8") as f:
        reader = csv.reader(f)
        rows = list(reader)
        for row in rows[1:]:
            st.markdown(f"**ğŸ—“ {row[1]}** â€“ `{row[0]}`  
{row[2]}")
except:
    st.info("ChÆ°a cÃ³ cuá»™c há»p nÃ o Ä‘Æ°á»£c lÆ°u.")
